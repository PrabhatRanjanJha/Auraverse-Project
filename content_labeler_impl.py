"""
content_labeler_impl.py

Text/document topic extractor using YAKE (offline keyword extraction).
Returns a single short folder-safe token.

Dependencies:
    pip install yake PyPDF2 python-docx python-pptx pandas beautifulsoup4 lxml
"""
from pathlib import Path
from typing import Optional
import re

def extract_text(path: str, ext: str) -> str:
    try:
        if ext == "txt":
            return Path(path).read_text(encoding="utf-8", errors="ignore")
        elif ext == "pdf":
            from PyPDF2 import PdfReader
            reader = PdfReader(path)
            return "\n".join((p.extract_text() or "") for p in reader.pages)
        elif ext in ("doc", "docx"):
            import docx
            doc = docx.Document(path)
            return "\n".join(p.text for p in doc.paragraphs)
        elif ext == "pptx":
            from pptx import Presentation
            prs = Presentation(path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        parts.append(shape.text)
            return "\n".join(parts)
        elif ext in {"csv", "xls", "xlsx"}:
            import pandas as pd
            try:
                df = pd.read_csv(path)
            except Exception:
                df = pd.read_excel(path)
            return df.astype(str).to_string()
        elif ext in {"html", "htm"}:
            from bs4 import BeautifulSoup
            html = Path(path).read_text(encoding="utf-8", errors="ignore")
            soup = BeautifulSoup(html, "lxml")
            return soup.get_text(separator="\n")
        elif ext == "md":
            return Path(path).read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        print(f"[content_labeler_impl] extraction error: {e}")
    return ""

def clean_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()

def infer_topic_yake(text: str) -> str:
    try:
        import yake
    except Exception:
        # fallback: simple keyword
        words = re.findall(r"[A-Za-z]{4,}", text.lower())
        from collections import Counter
        return Counter(words).most_common(1)[0][0] if words else "unknown"

    text = clean_text(text)
    if len(text) < 20:
        return "unknown"

    kw_extractor = yake.KeywordExtractor(lan="en", n=1, top=5, dedupLim=0.9)
    keywords = kw_extractor.extract_keywords(text)
    if not keywords:
        return "unknown"
    best = keywords[0][0].lower()
    best = re.sub(r"[^a-z0-9]+", "_", best).strip("_")
    return best or "unknown"

def label_from_content(path: str, file_type: Optional[str] = None) -> str:
    ext = (file_type or Path(path).suffix.lstrip(".")).lower()
    text = extract_text(path, ext)
    if text and text.strip():
        return infer_topic_yake(text)
    # fallback filename
    name = Path(path).stem.lower()
    name = re.sub(r"[^a-z0-9]+", "_", name).strip("_")
    return name or "unknown"
